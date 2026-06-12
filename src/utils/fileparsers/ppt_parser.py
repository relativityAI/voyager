"""
ppt_parser.py — Extract text (and optionally slide images) from PowerPoint files.

Supported formats:
    .pptx, .ppsx            → handled natively by python-pptx
    .ppt,  .pps  (legacy)   → converted via LibreOffice then handled as above

Requires:
    pip install python-pptx Pillow
    # For legacy .ppt files: LibreOffice must be installed system-wide
    #   Ubuntu/Debian: sudo apt install libreoffice
"""

from __future__ import annotations

import io
import shutil
import subprocess
import tempfile
from pathlib import Path

from loguru import logger

from src.utils.fileparsers.base import BaseParser, ParseResult

# Extensions supported natively by python-pptx
_PPTX_NATIVE = {".pptx", ".ppsx", ".pptm", ".ppsm"}
# Legacy binary formats — need LibreOffice conversion
_PPTX_LEGACY = {".ppt", ".pps"}


class PPTParser(BaseParser):
    """Parse PowerPoint presentations into text.

    Text is extracted from:
    - All text frames (titles, bodies, content placeholders)
    - Speaker notes
    - Table cells

    Args:
        slide_separator: String inserted between slides (default: ``"\\n\\n--- Slide {n} ---\\n\\n"``)
    """

    label = "PPTParser"

    def __init__(self, slide_separator: str | None = None) -> None:
        # None → use default pattern below
        self._custom_separator = slide_separator

    def _slide_sep(self, slide_num: int) -> str:
        if self._custom_separator is not None:
            return self._custom_separator
        return f"\n\n--- Slide {slide_num} ---\n\n"

    # ------------------------------------------------------------------
    def parse(
        self,
        source: str,
        *,
        extract_images: bool = False,
    ) -> ParseResult:
        """Extract text (and optionally slide thumbnails) from a presentation.

        Args:
            source:         Absolute path to the presentation file.
            extract_images: When ``True``, render each slide as a PNG thumbnail.
                            Requires ``Pillow`` and, for thumbnail rendering,
                            ``LibreOffice`` (``soffice``) or a similar tool.

        Returns:
            :class:`~src.utils.fileparsers.base.ParseResult`
        """
        try:
            from pptx import Presentation  # noqa: PLC0415
        except ImportError as exc:
            raise ImportError(
                "python-pptx is required: pip install python-pptx"
            ) from exc

        path = Path(source)
        suffix = path.suffix.lower()

        if suffix not in _PPTX_NATIVE | _PPTX_LEGACY:
            raise ValueError(
                f"Unsupported extension '{suffix}'. "
                f"Supported: {sorted(_PPTX_NATIVE | _PPTX_LEGACY)}"
            )

        if not path.exists():
            raise FileNotFoundError(f"Presentation not found: {source}")

        # --- Convert legacy .ppt → .pptx if needed ---
        working_path = path
        _tmp_dir = None
        if suffix in _PPTX_LEGACY:
            working_path, _tmp_dir = self._convert_legacy(path)

        try:
            prs = Presentation(working_path)
            texts: list[str] = []
            images: list = []
            total_slides = len(prs.slides)

            for slide_num, slide in enumerate(prs.slides, start=1):
                slide_text_parts: list[str] = []

                # --- Text frames ---
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            para_text = " ".join(
                                run.text for run in para.runs if run.text
                            )
                            if para_text.strip():
                                slide_text_parts.append(para_text)

                    # --- Table cells ---
                    if shape.has_table:
                        for row in shape.table.rows:
                            row_text = " | ".join(
                                cell.text_frame.text
                                for cell in row.cells
                                if cell.text_frame.text.strip()
                            )
                            if row_text.strip():
                                slide_text_parts.append(row_text)

                # --- Speaker notes ---
                if slide.has_notes_slide:
                    notes_tf = slide.notes_slide.notes_text_frame
                    notes_text = notes_tf.text.strip() if notes_tf else ""
                    if notes_text:
                        slide_text_parts.append(f"[Notes] {notes_text}")

                slide_text = "\n".join(slide_text_parts)
                texts.append(self._slide_sep(slide_num) + slide_text)

                # --- Slide thumbnail ---
                if extract_images:
                    img = self._render_slide_thumbnail(slide, slide_num)
                    if img is not None:
                        images.append(img)

        finally:
            if _tmp_dir is not None:
                shutil.rmtree(_tmp_dir, ignore_errors=True)

        return ParseResult(
            text="\n".join(texts).strip(),
            images=images,
            metadata={
                "total_slides": total_slides,
                "source_format": suffix,
            },
            source=str(source),
        )

    # ------------------------------------------------------------------
    # Helpers
    # ------------------------------------------------------------------

    @staticmethod
    def _convert_legacy(path: Path) -> tuple[Path, str]:
        """Use LibreOffice to convert legacy .ppt/.pps → .pptx.

        Returns the path to the converted .pptx and the temp dir that holds it.
        Caller is responsible for cleaning up the temp dir.
        """
        soffice = shutil.which("soffice") or shutil.which("libreoffice")
        if soffice is None:
            raise RuntimeError(
                "LibreOffice (soffice) not found in PATH. "
                "Install it with:  sudo apt install libreoffice  "
                "to handle legacy .ppt/.pps files."
            )

        tmp_dir = tempfile.mkdtemp(prefix="voyager_ppt_")
        try:
            subprocess.run(
                [
                    soffice,
                    "--headless",
                    "--convert-to", "pptx",
                    "--outdir", tmp_dir,
                    str(path),
                ],
                check=True,
                capture_output=True,
                timeout=120,
            )
        except subprocess.CalledProcessError as exc:
            shutil.rmtree(tmp_dir, ignore_errors=True)
            raise RuntimeError(
                f"LibreOffice failed to convert {path}: {exc.stderr.decode()}"
            ) from exc

        converted = list(Path(tmp_dir).glob("*.pptx"))
        if not converted:
            shutil.rmtree(tmp_dir, ignore_errors=True)
            raise RuntimeError(
                f"LibreOffice did not produce a .pptx file in {tmp_dir}"
            )

        return converted[0], tmp_dir

    @staticmethod
    def _render_slide_thumbnail(slide, slide_num: int):
        """Attempt to render a slide as a PIL Image (best-effort)."""
        try:
            from PIL import Image, ImageDraw, ImageFont  # noqa: PLC0415

            # python-pptx doesn't render slides natively — we create a
            # simple placeholder image containing the slide text.
            # Full rendering requires external tools (LibreOffice, comtypes on Win).
            width, height = 800, 600
            img = Image.new("RGB", (width, height), color=(255, 255, 255))
            draw = ImageDraw.Draw(img)

            # Collect all text from the slide
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    texts.append(shape.text_frame.text.strip())
            combined = "\n".join(t for t in texts if t)

            try:
                font = ImageFont.truetype("DejaVuSans.ttf", 20)
            except Exception:
                font = ImageFont.load_default()

            draw.text((20, 20), f"Slide {slide_num}", fill=(50, 50, 50), font=font)
            draw.text((20, 60), combined[:800], fill=(0, 0, 0), font=font)
            return img

        except Exception as err:
            logger.warning(f"[PPTParser] Could not render slide {slide_num}: {err}")
            return None
